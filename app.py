"""
File Analyser — Production Chainlit App
Supports any file type: documents, spreadsheets, code, images, data files, and more.
"""

import chainlit as cl
import pandas as pd
import docx
from pptx import Presentation
import os
import json
import base64
from pathlib import Path
from bs4 import BeautifulSoup
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.output_parsers import StrOutputParser
from langchain_core.messages import HumanMessage
from langchain_community.document_loaders import PyMuPDFLoader, BSHTMLLoader
import chardet

# ═══════════════════════════════════════════════════════════════════════════════
# CONFIGURATION
# ═══════════════════════════════════════════════════════════════════════════════

os.environ["CHAINLIT_TELEMETRY_ENABLED"] = "false"

MAX_CONTEXT_CHARS      = 14_000   # chars sent to LLM for Q&A
MAX_TRANSLATION_CHARS  = 12_000   # chars sent for translation
MAX_FILE_SIZE_MB       = 50       # upload limit
MAX_HISTORY_PAIRS      = 5        # Q&A conversation history pairs kept

# App states stored in user session
STATE_IDLE    = "idle"     # waiting for file upload
STATE_LOADED  = "loaded"   # file loaded, waiting for action button
STATE_ASKING  = "asking"   # Q&A mode, chat input active

LANGUAGES = [
    "Arabic", "Bengali", "Chinese (Simplified)", "Chinese (Traditional)",
    "Dutch", "English", "French", "German", "Greek", "Hindi",
    "Indonesian", "Italian", "Japanese", "Korean", "Malay",
    "Polish", "Portuguese", "Russian", "Spanish", "Swedish",
    "Tamil", "Thai", "Turkish", "Ukrainian", "Vietnamese",
]

# Extension → processing category
EXT_CATEGORY: dict[str, str] = {
    # Spreadsheets
    ".csv": "spreadsheet", ".xlsx": "spreadsheet", ".xls": "spreadsheet",
    ".ods": "spreadsheet", ".tsv": "spreadsheet",
    # Word documents
    ".docx": "document", ".doc": "document",
    # Presentations
    ".pptx": "presentation", ".ppt": "presentation",
    # PDF
    ".pdf": "pdf",
    # Web
    ".html": "html", ".htm": "html",
    # Markup / data
    ".xml": "xml", ".xsd": "xml", ".xsl": "xml", ".svg": "xml",
    ".json": "json", ".jsonl": "json", ".geojson": "json",
    # Images (analysed via Gemini Vision — no extra deps needed)
    ".png": "image", ".jpg": "image", ".jpeg": "image",
    ".webp": "image", ".gif": "image", ".bmp": "image", ".tiff": "image",
    # Plain text and all code / config files → read as text
    ".txt": "text", ".md": "text", ".markdown": "text", ".rst": "text",
    ".log": "text", ".yaml": "text", ".yml": "text", ".toml": "text",
    ".ini": "text", ".cfg": "text", ".conf": "text", ".env": "text",
    ".py": "text",  ".pyw": "text", ".js": "text",  ".ts": "text",
    ".jsx": "text", ".tsx": "text", ".java": "text", ".c": "text",
    ".cpp": "text", ".cc": "text",  ".h": "text",   ".hpp": "text",
    ".cs": "text",  ".go": "text",  ".rs": "text",  ".rb": "text",
    ".php": "text", ".swift": "text", ".kt": "text", ".scala": "text",
    ".sql": "text", ".sh": "text",  ".bash": "text", ".zsh": "text",
    ".ps1": "text", ".bat": "text", ".cmd": "text",
    ".r": "text",   ".m": "text",   ".lua": "text",  ".pl": "text",
    ".dockerfile": "text", ".makefile": "text",
}

# ═══════════════════════════════════════════════════════════════════════════════
# LLM INITIALIZATION
# ═══════════════════════════════════════════════════════════════════════════════

def _init_llm():
    key = os.environ.get("GOOGLE_API_KEY", "").strip()
    if not key:
        return None
    try:
        return ChatGoogleGenerativeAI(
            model="gemini-2.5-flash",
            google_api_key=key,
            temperature=0.1,
            streaming=True,
        )
    except Exception as exc:
        print(f"[LLM] Init failed: {exc}")
        return None


llm = _init_llm()

# ═══════════════════════════════════════════════════════════════════════════════
# FILE UTILITIES
# ═══════════════════════════════════════════════════════════════════════════════

def get_category(file) -> str:
    """Determine processing category from file extension."""
    ext = Path(file.name).suffix.lower()
    return EXT_CATEGORY.get(ext, "text")  # unknown extension → try as text


def safe_decode(path: str) -> str:
    """Read a file as text, auto-detecting its character encoding."""
    with open(path, "rb") as fh:
        raw = fh.read()
    enc = chardet.detect(raw).get("encoding") or "utf-8"
    try:
        return raw.decode(enc)
    except Exception:
        return raw.decode("utf-8", errors="replace")


def human_size(n: int) -> str:
    if n < 1_024:      return f"{n} B"
    if n < 1_048_576:  return f"{n / 1_024:.1f} KB"
    return f"{n / 1_048_576:.1f} MB"


def load_dataframe(file) -> pd.DataFrame:
    """Load any spreadsheet format into a Pandas DataFrame."""
    ext  = Path(file.name).suffix.lower()
    path = file.path
    if ext in {".csv", ".tsv"}:
        sep = "\t" if ext == ".tsv" else ","
        with open(path, "rb") as fh:
            enc = chardet.detect(fh.read(8_192)).get("encoding") or "utf-8"
        return pd.read_csv(path, encoding=enc, sep=sep)
    if ext in {".xlsx", ".xlsm"}:
        return pd.read_excel(path, engine="openpyxl")
    if ext == ".xls":
        return pd.read_excel(path, engine="xlrd")
    if ext == ".ods":
        return pd.read_excel(path, engine="odf")
    return pd.read_excel(path)


async def extract_content(file) -> tuple:
    """
    Parse any uploaded file into usable content.

    Returns
    -------
    content  : str | pd.DataFrame | dict   (dict for images: {mime, b64})
    category : str                         ("text" | "spreadsheet" | "image")
    stats    : dict                        (metadata for UI display)
    """
    path     = file.path
    ext      = Path(file.name).suffix.lower()
    size     = os.path.getsize(path)
    category = get_category(file)
    stats    = {"size": human_size(size), "ext": ext or "?"}

    try:
        # ── SPREADSHEET ────────────────────────────────────────────────────
        if category == "spreadsheet":
            df = load_dataframe(file)
            stats.update(rows=df.shape[0], cols=df.shape[1], category="spreadsheet")
            return df, "spreadsheet", stats

        # ── IMAGE → Gemini Vision ──────────────────────────────────────────
        elif category == "image":
            mime = file.type or f"image/{ext.lstrip('.').replace('jpg', 'jpeg')}"
            with open(path, "rb") as fh:
                b64 = base64.b64encode(fh.read()).decode()
            stats["category"] = "image"
            return {"mime": mime, "b64": b64}, "image", stats

        # ── WORD DOCUMENT ──────────────────────────────────────────────────
        elif category == "document":
            d     = docx.Document(path)
            parts = [p.text for p in d.paragraphs if p.text.strip()]
            # Also extract table content
            for table in d.tables:
                for row in table.rows:
                    row_txt = " | ".join(
                        c.text.strip() for c in row.cells if c.text.strip()
                    )
                    if row_txt:
                        parts.append(row_txt)
            text = "\n".join(parts)
            stats.update(chars=len(text), category="document")
            return text, "text", stats

        # ── PDF ────────────────────────────────────────────────────────────
        elif category == "pdf":
            loader = PyMuPDFLoader(path)
            pages  = loader.load()
            text   = "\n\n".join(p.page_content for p in pages)
            stats.update(pages=len(pages), chars=len(text), category="pdf")
            return text, "text", stats

        # ── POWERPOINT ─────────────────────────────────────────────────────
        elif category == "presentation":
            prs   = Presentation(path)
            parts = []
            for i, slide in enumerate(prs.slides, 1):
                parts.append(f"\n── Slide {i} ──")
                for shape in slide.shapes:
                    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            if para.text.strip():
                                parts.append(para.text)
            text = "\n".join(parts)
            stats.update(slides=len(prs.slides), chars=len(text), category="presentation")
            return text, "text", stats

        # ── HTML ───────────────────────────────────────────────────────────
        elif category == "html":
            loader = BSHTMLLoader(path, open_encoding="utf-8")
            docs   = loader.load()
            text   = "\n".join(d.page_content for d in docs)
            stats.update(chars=len(text), category="html")
            return text, "text", stats

        # ── XML / SVG ──────────────────────────────────────────────────────
        elif category == "xml":
            raw  = safe_decode(path)
            soup = BeautifulSoup(raw, "xml")
            text = soup.get_text(separator="\n", strip=True)
            stats.update(chars=len(text), category="xml")
            return text, "text", stats

        # ── JSON / JSONL ───────────────────────────────────────────────────
        elif category == "json":
            raw = safe_decode(path)
            if ext == ".jsonl":
                lines = [
                    json.loads(ln)
                    for ln in raw.strip().splitlines()
                    if ln.strip()
                ]
                text = json.dumps(lines[:200], indent=2, ensure_ascii=False)
            else:
                text = json.dumps(json.loads(raw), indent=2, ensure_ascii=False)
            stats.update(chars=len(text), category="json")
            return text, "text", stats

        # ── PLAIN TEXT / CODE / ANYTHING ELSE ─────────────────────────────
        else:
            try:
                text = safe_decode(path)
                stats.update(chars=len(text), category="text")
                return text, "text", stats
            except Exception:
                return None, "binary", stats

    except Exception as exc:
        raise RuntimeError(f"Could not parse `{file.name}`: {exc}") from exc


# ═══════════════════════════════════════════════════════════════════════════════
# UI HELPERS
# ═══════════════════════════════════════════════════════════════════════════════

def _main_actions(include_new_file: bool = False) -> list:
    actions = [
        cl.Action(name="translate",     label="🌐  Translate",     payload={}),
        cl.Action(name="ask_questions", label="💬  Ask Questions", payload={}),
    ]
    if include_new_file:
        actions.append(cl.Action(name="new_file", label="📁  New File", payload={}))
    return actions


async def show_main_menu(include_new_file: bool = False):
    cl.user_session.set("state", STATE_LOADED)
    await cl.Message(
        content="**What would you like to do?**",
        actions=_main_actions(include_new_file),
    ).send()


def _stats_line(file_name: str, stats: dict) -> str:
    """Build the ✅ summary line shown after a file loads."""
    parts = [f"**{file_name}**", f"`{stats['size']}`"]
    cat   = stats.get("category", "")
    if cat == "spreadsheet":
        parts.append(f"`{stats['rows']:,} rows × {stats['cols']} columns`")
    if "pages" in stats:
        parts.append(f"`{stats['pages']} pages`")
    if "slides" in stats:
        parts.append(f"`{stats['slides']} slides`")
    if "chars" in stats:
        parts.append(f"`{stats['chars']:,} characters`")
    if cat == "image":
        parts.append("🖼️ image")
    return "✅  " + "  ·  ".join(parts)


# ═══════════════════════════════════════════════════════════════════════════════
# CORE UPLOAD FLOW
# ═══════════════════════════════════════════════════════════════════════════════

@cl.on_chat_start
async def start():
    if not llm:
        await cl.Message(
            content=(
                "## ⚠️ Missing API Key\n"
                "`GOOGLE_API_KEY` is not set or is invalid.\n\n"
                "Set the environment variable and restart the app."
            )
        ).send()
        return

    await cl.Message(
        content=(
            "# 📂 File Analyser\n"
            "Upload **any file** — I'll extract its content so you can "
            "**ask questions** or **translate** it.\n\n"
            "> **Supported:** Documents · PDFs · Spreadsheets · Presentations · "
            "Code · JSON · XML · HTML · Images · and more"
        )
    ).send()

    await _handle_upload()


async def _handle_upload():
    """Ask for a file upload, parse it, persist in session, show main menu."""
    cl.user_session.set("state", STATE_IDLE)

    files = await cl.AskFileMessage(
        content="Upload your file:",
        accept=["*/*"],
        max_size_mb=MAX_FILE_SIZE_MB,
        timeout=600,
    ).send()

    if not files:
        await cl.Message(content="⚠️ No file received. Please try again.").send()
        return

    file       = files[0]
    status_msg = cl.Message(content=f"⏳ Loading `{file.name}`…")
    await status_msg.send()

    try:
        content, category, stats = await extract_content(file)
    except RuntimeError as exc:
        status_msg.content = f"❌ {exc}"
        await status_msg.update()
        return

    if content is None:
        status_msg.content = (
            f"❌ `{file.name}` is a binary file and cannot be read as text.\n"
            "Please try a different file."
        )
        await status_msg.update()
        return

    # Persist everything in the session
    cl.user_session.set("content",      content)
    cl.user_session.set("category",     category)
    cl.user_session.set("file_name",    file.name)
    cl.user_session.set("stats",        stats)
    cl.user_session.set("chat_history", [])

    status_msg.content = _stats_line(file.name, stats)
    await status_msg.update()

    await show_main_menu(include_new_file=False)


# ═══════════════════════════════════════════════════════════════════════════════
# ACTION CALLBACKS
# ═══════════════════════════════════════════════════════════════════════════════

@cl.action_callback("new_file")
async def on_new_file(_action: cl.Action):
    cl.user_session.set("content",      None)
    cl.user_session.set("chat_history", [])
    await _handle_upload()


@cl.action_callback("back_to_menu")
async def on_back_to_menu(_action: cl.Action):
    await show_main_menu(include_new_file=True)


@cl.action_callback("ask_questions")
async def on_ask_questions(_action: cl.Action):
    cl.user_session.set("state", STATE_ASKING)
    file_name = cl.user_session.get("file_name", "the file")
    category  = cl.user_session.get("category", "")

    hints = {
        "spreadsheet": "\n\n> 💡 Ask for summaries, comparisons, averages, trends, or specific values.",
        "image":       "\n\n> 💡 Ask me to describe the image, read any text in it, or identify objects.",
        "text":        "\n\n> 💡 Ask for summaries, key points, specific sections, or explanations.",
    }
    hint = hints.get(category, "")

    await cl.Message(
        content=(
            f"💬 **Ask Questions mode** — `{file_name}`\n"
            f"Type your questions below. I'll answer using only the file content.{hint}"
        ),
        actions=[
            cl.Action(name="back_to_menu", label="← Back to Menu", payload={}),
            cl.Action(name="new_file",     label="📁 New File",     payload={}),
        ],
    ).send()


@cl.action_callback("translate")
async def on_translate(_action: cl.Action):
    category = cl.user_session.get("category", "")

    # Unsupported types for translation
    if category == "image":
        await cl.Message(
            content=(
                "🖼️ **Image translation is not supported.**\n"
                "Switch to **Ask Questions** to describe or extract text from the image."
            )
        ).send()
        await show_main_menu(include_new_file=True)
        return

    if category == "spreadsheet":
        await cl.Message(
            content=(
                "📊 **Spreadsheet translation is not supported.**\n"
                "Switch to **Ask Questions** to query, summarise, or describe the data."
            )
        ).send()
        await show_main_menu(include_new_file=True)
        return

    # Show language picker
    lang_actions = [
        cl.Action(name=f"lang_{i}", label=lang, payload={"lang": lang})
        for i, lang in enumerate(LANGUAGES)
    ]
    lang_res = await cl.AskActionMessage(
        content="🌍 **Select target language:**",
        actions=lang_actions,
        timeout=120,
    ).send()

    if not lang_res:
        await cl.Message(content="Translation cancelled.").send()
        await show_main_menu(include_new_file=True)
        return

    target_lang = lang_res.get("payload", {}).get("lang", "English")
    await _run_translation(target_lang)


# Register one callback per language (lang_0 … lang_N)
for _idx, _language in enumerate(LANGUAGES):
    def _make_cb(lang: str):
        async def _cb(_action: cl.Action):
            await _run_translation(lang)
        return _cb
    cl.action_callback(f"lang_{_idx}")(_make_cb(_language))


async def _run_translation(target_lang: str):
    """Stream-translate the loaded document to the chosen language."""
    content   = cl.user_session.get("content")
    file_name = cl.user_session.get("file_name", "document")
    text      = content if isinstance(content, str) else ""

    if len(text) > MAX_TRANSLATION_CHARS:
        await cl.Message(
            content=(
                f"⚠️ Document is large ({len(text):,} chars). "
                f"Translating the first **{MAX_TRANSLATION_CHARS:,} characters**."
            )
        ).send()
        text = text[:MAX_TRANSLATION_CHARS]

    msg = cl.Message(content=f"🌐 **Translation → {target_lang}**\n\n")
    await msg.send()
    try:
        prompt = ChatPromptTemplate.from_template(
            "Translate the following text to {target_lang} naturally and accurately.\n"
            "Preserve all structure: headings, bullet points, tables, numbered lists, and code blocks.\n"
            "Output ONLY the translated text — no preamble, no explanation.\n\n"
            "{text}"
        )
        chain = prompt | llm | StrOutputParser()
        async for chunk in chain.astream({"target_lang": target_lang, "text": text}):
            await msg.stream_token(chunk)
    except Exception as exc:
        await msg.stream_token(f"\n\n❌ Translation failed: {exc}")

    await show_main_menu(include_new_file=True)
    """Stream-translate the loaded document to the chosen language."""
    content   = cl.user_session.get("content")
    file_name = cl.user_session.get("file_name", "document")
    text      = content if isinstance(content, str) else ""

    if len(text) > MAX_TRANSLATION_CHARS:
        await cl.Message(
            content=(
                f"⚠️ Document is large ({len(text):,} chars). "
                f"Translating the first **{MAX_TRANSLATION_CHARS:,} characters**."
            )
        ).send()
        text = text[:MAX_TRANSLATION_CHARS]

    async with cl.Message(content=f"🌐 **Translation → {target_lang}**\n\n") as msg:
        try:
            prompt = ChatPromptTemplate.from_template(
                "Translate the following text to {target_lang} naturally and accurately.\n"
                "Preserve all structure: headings, bullet points, tables, numbered lists, and code blocks.\n"
                "Output ONLY the translated text — no preamble, no explanation.\n\n"
                "{text}"
            )
            chain = prompt | llm | StrOutputParser()
            async for chunk in chain.astream({"target_lang": target_lang, "text": text}):
                await msg.stream_token(chunk)
        except Exception as exc:
            await msg.stream_token(f"\n\n❌ Translation failed: {exc}")

    await show_main_menu(include_new_file=True)


# ═══════════════════════════════════════════════════════════════════════════════
# Q&A MESSAGE HANDLER
# ═══════════════════════════════════════════════════════════════════════════════

@cl.on_message
async def on_message(message: cl.Message):
    state     = cl.user_session.get("state", STATE_IDLE)
    content   = cl.user_session.get("content")
    category  = cl.user_session.get("category", "text")
    file_name = cl.user_session.get("file_name", "the file")
    question  = message.content.strip()

    if state != STATE_ASKING:
        await cl.Message(
            content="👆 Please use the **buttons above** to choose an action first."
        ).send()
        return

    if not content:
        await cl.Message(content="⚠️ No file loaded. Please upload a file first.").send()
        return

    if not question:
        return

    reply = cl.Message(content="")
    await reply.send()

    try:
        # ── IMAGE → Gemini Vision ──────────────────────────────────────
        if category == "image":
            data_url = f"data:{content['mime']};base64,{content['b64']}"
            resp = await llm.ainvoke([
                HumanMessage(content=[
                    {"type": "image_url", "image_url": {"url": data_url}},
                    {"type": "text",      "text": question},
                ])
            ])
            await reply.stream_token(resp.content)

        # ── SPREADSHEET ────────────────────────────────────────────────
        elif category == "spreadsheet":
            df: pd.DataFrame = content
            num_cols = df.select_dtypes(include="number").columns.tolist()
            cat_cols = df.select_dtypes(exclude="number").columns.tolist()
            summary  = (
                f"File: {file_name}\n"
                f"Shape: {df.shape[0]:,} rows × {df.shape[1]} columns\n"
                f"All columns: {list(df.columns)}\n"
                f"Numeric columns: {num_cols}\n"
                f"Categorical columns: {cat_cols}\n\n"
                f"Data types:\n{df.dtypes.to_string()}\n\n"
                f"First 5 rows:\n{df.head(5).to_string()}\n\n"
                f"Descriptive statistics:\n{df.describe(include='all').to_string()}"
            )
            prompt = ChatPromptTemplate.from_template(
                "You are an expert data analyst. Answer using ONLY the dataset information below.\n"
                "Show your reasoning step-by-step when calculations are involved.\n"
                "If the data is insufficient to answer, say so clearly.\n\n"
                "Dataset — {file_name}:\n{summary}\n\n"
                "Question: {question}"
            )
            chain = prompt | llm | StrOutputParser()
            async for chunk in chain.astream(
                {"file_name": file_name, "summary": summary, "question": question}
            ):
                await reply.stream_token(chunk)

        # ── TEXT / DOCUMENTS / CODE ────────────────────────────────────
        else:
            history: list[dict] = cl.user_session.get("chat_history", [])
            history_str = ""
            if history:
                recent = history[-MAX_HISTORY_PAIRS:]
                history_str = "Previous conversation:\n" + "\n".join(
                    f"Q: {h['q']}\nA: {h['a']}" for h in recent
                ) + "\n\n"

            text      = str(content)
            truncated = len(text) > MAX_CONTEXT_CHARS
            ctx       = text[:MAX_CONTEXT_CHARS]
            trunc_note = (
                f" [first {MAX_CONTEXT_CHARS:,} of {len(text):,} chars]"
                if truncated else ""
            )

            prompt = ChatPromptTemplate.from_template(
                "You are an expert document analyst. Answer questions based SOLELY on the "
                "document content provided below. Be thorough, specific, and well-structured.\n"
                "Quote relevant passages when it adds clarity.\n"
                "If the answer is not found in the document, say so explicitly — do not guess.\n\n"
                "{history}"
                "Document — {file_name}{trunc_note}:\n"
                "```\n{context}\n```\n\n"
                "Question: {question}"
            )
            chain       = prompt | llm | StrOutputParser()
            full_answer = ""
            async for chunk in chain.astream({
                "history":    history_str,
                "file_name":  file_name,
                "trunc_note": trunc_note,
                "context":    ctx,
                "question":   question,
            }):
                full_answer += chunk
                await reply.stream_token(chunk)

            history.append({"q": question, "a": full_answer})
            cl.user_session.set("chat_history", history)

    except Exception as exc:
        await reply.stream_token(f"\n\n❌ Error: {exc}")
